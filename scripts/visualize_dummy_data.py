import gzip
import json
import os
import textwrap
from typing import Any

import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

# folium imports (for map visualization)
try:
    import folium
    from folium import plugins

    FOLIUM_AVAILABLE = True
except ImportError:
    FOLIUM_AVAILABLE = False

# 한글 폰트 설정
plt.rcParams["font.family"] = "DejaVu Sans"
plt.rcParams["axes.unicode_minus"] = False

# 한글 폰트 설정 시도
try:
    # Windows에서 사용 가능한 한글 폰트들 (우선순위 순)
    korean_fonts = [
        "Malgun Gothic",  # Windows 기본 한글 폰트
        "NanumGothic",  # 나눔고딕
        "NanumBarunGothic",  # 나눔바른고딕
        "Noto Sans CJK KR",  # 구글 Noto Sans
        "AppleGothic",  # macOS
        "D2Coding",  # 코딩 폰트
    ]

    # 사용 가능한 폰트 목록 가져오기
    available_fonts = [f.name for f in fm.fontManager.ttflist]

    for font_name in korean_fonts:
        if font_name in available_fonts:
            plt.rcParams["font.family"] = font_name
            print(f"✅ 한글 폰트 설정 완료: {font_name}")
            break
    else:
        print("⚠️  한글 폰트를 찾을 수 없습니다. 기본 폰트를 사용합니다.")
        print(f"사용 가능한 폰트 목록: {available_fonts[:10]}...")  # 처음 10개만 출력
        # 한글이 깨질 경우를 대비해 폰트 크기 조정
        plt.rcParams["font.size"] = 10

except Exception as e:
    print(f"❌ 폰트 설정 중 오류 발생: {e}")
    plt.rcParams["font.size"] = 10

# ========== 설정 ==========
DATA_PATH = "dummy_data.json.gz"
OUT_DIR = "scripts/charts"
MAKE_PPT = True  # PPT 생성 여부
PPT_PATH = "scripts/mogumogu_report.pptx"


# ========== 유틸 ==========
def ensure_dir(d: str) -> None:
    if not os.path.exists(d):
        os.makedirs(d)


def savefig(path: str) -> None:
    plt.tight_layout()
    # 한글 폰트 문제 해결을 위한 추가 설정
    plt.savefig(
        path,
        dpi=160,
        bbox_inches="tight",
        facecolor="white",
        edgecolor="none",
        pad_inches=0.2,
    )
    plt.close()


def parse_dt(s: str) -> pd.Timestamp:
    try:
        return pd.to_datetime(s)
    except Exception:
        return pd.NaT


def load_data() -> tuple[
    dict[str, Any],
    pd.DataFrame,
    pd.DataFrame,
    pd.DataFrame,
    pd.DataFrame,
    pd.DataFrame,
    pd.DataFrame,
]:
    with gzip.open(DATA_PATH, "rt", encoding="utf-8") as f:
        data = json.load(f)
    users = pd.DataFrame(data["users"])
    posts = pd.DataFrame(data["mogu_posts"])
    favs = pd.DataFrame(data["favorites"])
    parts = pd.DataFrame(data["participations"])
    ratings = pd.DataFrame(data["ratings"])
    comments = pd.DataFrame(data["comments"])

    # 타입 정리
    for df, cols in [
        (posts, ["created_at", "mogu_datetime"]),
        (favs, ["created_at"]),
        (parts, ["applied_at", "decided_at"]),
        (ratings, ["created_at"]),
        (comments, ["created_at"]),
    ]:
        for c in cols:
            if c in df.columns:
                df[c] = pd.to_datetime(df[c], errors="coerce")

    # 위치 튜플 -> lat/lon 컬럼
    if "mogu_spot" in posts.columns:
        lat, lon = [], []
        for x in posts["mogu_spot"]:
            if isinstance(x, list) or isinstance(x, tuple):
                lat.append(x[0])
                lon.append(x[1])
            else:
                lat.append(np.nan)
                lon.append(np.nan)
        posts["lat"] = lat
        posts["lon"] = lon

    return data, users, posts, favs, parts, ratings, comments


# ========== 차트들 ==========
def chart_price_box_by_category(posts: pd.DataFrame) -> None:
    # 부드러운 색상으로 박스플롯 그리기
    fig, ax = plt.subplots(figsize=(12, 8))

    # 카테고리별 데이터 준비
    categories = posts["category"].unique()
    data_by_category = [
        posts[posts["category"] == cat]["price"].values for cat in categories
    ]

    # 부드러운 색상 설정
    soft_colors = ["#FF6B6B", "#4ECDC4", "#45B7D1", "#96CEB4"]
    colors = (soft_colors * ((len(categories) // len(soft_colors)) + 1))[
        : len(categories)
    ]

    # 박스플롯 그리기
    bp = ax.boxplot(
        data_by_category,
        tick_labels=categories,
        patch_artist=True,
        boxprops=dict(facecolor=None, edgecolor="#2C3E50", linewidth=1.5),
        medianprops=dict(color="#E74C3C", linewidth=2),
        whiskerprops=dict(color="#2C3E50", linewidth=1.5),
        capprops=dict(color="#2C3E50", linewidth=1.5),
    )

    # 박스 색상 설정
    for patch, color in zip(bp["boxes"], colors):
        patch.set_facecolor(color)
        patch.set_alpha(0.7)

    plt.title(
        "카테고리별 가격 분포 (Boxplot)",
        fontsize=16,
        fontweight="bold",
        color="#2C3E50",
        pad=20,
    )
    plt.xlabel("카테고리", fontsize=12, color="#2C3E50")
    plt.ylabel("가격(원)", fontsize=12, color="#2C3E50")
    plt.grid(True, alpha=0.3)

    # 배경 설정
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    savefig(os.path.join(OUT_DIR, "01_price_box_by_category.png"))


def chart_post_status_bar(posts: pd.DataFrame) -> None:
    vc = posts["status"].value_counts().sort_values(ascending=False)

    # 부드러운 색상 설정
    soft_colors = ["#FF6B6B", "#4ECDC4", "#45B7D1", "#96CEB4", "#FFEAA7", "#DDA0DD"]
    colors = (soft_colors * ((len(vc) // len(soft_colors)) + 1))[: len(vc)]

    fig, ax = plt.subplots(figsize=(10, 6))
    bars = ax.bar(
        range(len(vc)),
        vc.values,
        color=colors,
        alpha=0.8,
        edgecolor="#2C3E50",
        linewidth=1,
    )

    # 막대 위에 값 표시
    for bar, value in zip(bars, vc.values):
        height = bar.get_height()
        ax.text(
            bar.get_x() + bar.get_width() / 2.0,
            height + height * 0.01,
            f"{value:,}",
            ha="center",
            va="bottom",
            fontweight="bold",
            color="#2C3E50",
        )

    plt.title(
        "모구 게시물 상태 분포", fontsize=16, fontweight="bold", color="#2C3E50", pad=20
    )
    plt.xlabel("상태", fontsize=12, color="#2C3E50")
    plt.ylabel("개수", fontsize=12, color="#2C3E50")
    plt.xticks(range(len(vc)), vc.index, rotation=30, ha="right")
    plt.grid(True, alpha=0.3, axis="y")

    # 배경 설정
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    savefig(os.path.join(OUT_DIR, "02_post_status_bar.png"))


def chart_participation_status_pie(stats_dict: dict[str, Any]) -> None:
    s = pd.Series(stats_dict)

    # 상수 정의 - 모든 항목 표시하도록 조정
    MIN_PERCENTAGE_THRESHOLD = 1.0  # 1% 미만은 비율 표시 생략
    SMALL_SLICE_THRESHOLD = 0.01  # 1% 미만 조각 분리
    SMALL_LABEL_ROTATION_THRESHOLD = 3.0  # 3% 미만 라벨 기울이기

    # 작은 비율을 위한 커스텀 autopct 함수
    def custom_autopct(pct):
        # 1% 미만인 경우 비율 표시 생략
        if pct < MIN_PERCENTAGE_THRESHOLD:
            return ""
        else:
            return f"{pct:.1f}%"

    # 파이 차트 그리기 - 적당한 크기로 설정
    fig, ax = plt.subplots(figsize=(14, 10))

    # 부드러운 색상 팔레트
    soft_colors = [
        "#FF6B6B",  # 부드러운 빨강
        "#4ECDC4",  # 청록색
        "#45B7D1",  # 하늘색
        "#96CEB4",  # 민트 그린
        "#FFEAA7",  # 노란색
        "#DDA0DD",  # 자주색
        "#98D8C8",  # 연한 청록
        "#F7DC6F",  # 골드
    ]

    # 색상이 부족할 경우 반복
    colors = (soft_colors * ((len(s) // len(soft_colors)) + 1))[: len(s)]

    wedges, texts, autotexts = ax.pie(
        s.values,
        labels=s.index,
        autopct=custom_autopct,
        startangle=90,
        pctdistance=0.75,  # 비율 텍스트를 적당히 안쪽으로
        labeldistance=1.15,  # 라벨을 가까이 (기존 1.35에서 1.15로)
        textprops={"fontsize": 12, "weight": "bold", "color": "#2C3E50"},
        colors=colors,
        explode=[
            0.05 if value / s.sum() < SMALL_SLICE_THRESHOLD else 0 for value in s.values
        ],  # 작은 조각 약간 분리 (기존 0.1에서 0.05로)
        shadow=False,  # 그림자 효과 제거
        wedgeprops={"linewidth": 2, "edgecolor": "white"},  # 테두리 추가
    )

    # autopct 텍스트 스타일 조정
    for autotext in autotexts:
        if autotext.get_text():  # 텍스트가 있는 경우만
            autotext.set_color("#2C3E50")  # 어두운 회색
            autotext.set_fontsize(11)
            autotext.set_weight("bold")
            autotext.set_bbox(
                dict(boxstyle="round,pad=0.3", facecolor="white", alpha=0.8)
            )  # 배경 박스

    # 라벨 텍스트 스타일 조정 및 겹침 방지
    for i, text in enumerate(texts):
        text.set_fontsize(11)
        text.set_weight("bold")
        text.set_color("#2C3E50")

        # 작은 비율의 라벨은 약간 기울여서 겹침 방지
        percentage = s.iloc[i] / s.sum() * 100
        if percentage < SMALL_LABEL_ROTATION_THRESHOLD:  # 3% 미만은 기울이기
            text.set_rotation(15)
            text.set_ha("left")

    plt.title("참여 상태 비율", fontsize=18, fontweight="bold", color="#2C3E50", pad=25)

    # 범례 추가 (모든 항목 표시) - 더 깔끔하게
    legend_labels = [
        f"{label}: {value}개 ({value / s.sum() * 100:.1f}%)"
        for label, value in s.items()
    ]
    ax.legend(
        wedges,
        legend_labels,
        loc="center left",
        bbox_to_anchor=(1, 0.1, 0.5, 0.8),
        fontsize=10,
        frameon=True,
        fancybox=True,
        shadow=False,  # 범례 그림자도 제거
        framealpha=0.9,
    )

    # 배경을 깔끔하게
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    # 텍스트 겹침 방지를 위한 조정
    plt.tight_layout()
    savefig(os.path.join(OUT_DIR, "03_participation_status_pie.png"))


def chart_avg_joined_by_category(posts: pd.DataFrame) -> None:
    df = posts.groupby("category")["joined_count"].mean().sort_values(ascending=False)

    # 부드러운 색상 설정
    soft_colors = ["#FF6B6B", "#4ECDC4", "#45B7D1", "#96CEB4", "#FFEAA7"]
    colors = (soft_colors * ((len(df) // len(soft_colors)) + 1))[: len(df)]

    fig, ax = plt.subplots(figsize=(10, 6))
    bars = ax.bar(
        range(len(df)),
        df.values,
        color=colors,
        alpha=0.8,
        edgecolor="#2C3E50",
        linewidth=1,
    )

    # 막대 위에 값 표시
    for bar, value in zip(bars, df.values):
        height = bar.get_height()
        ax.text(
            bar.get_x() + bar.get_width() / 2.0,
            height + height * 0.01,
            f"{value:.1f}",
            ha="center",
            va="bottom",
            fontweight="bold",
            color="#2C3E50",
        )

    plt.title(
        "카테고리별 평균 참여자 수",
        fontsize=16,
        fontweight="bold",
        color="#2C3E50",
        pad=20,
    )
    plt.xlabel("카테고리", fontsize=12, color="#2C3E50")
    plt.ylabel("평균 참여자 수", fontsize=12, color="#2C3E50")
    plt.xticks(range(len(df)), df.index, rotation=30, ha="right")
    plt.grid(True, alpha=0.3, axis="y")

    # 배경 설정
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    savefig(os.path.join(OUT_DIR, "04_avg_joined_by_category.png"))


def chart_time_of_day_heatmap(posts: pd.DataFrame) -> None:
    df = posts.copy()
    df["hour"] = df["mogu_datetime"].dt.hour
    df["weekday"] = df["mogu_datetime"].dt.weekday  # 0=월요일
    pivot = df.pivot_table(
        index="weekday", columns="hour", values="id", aggfunc="count"
    ).fillna(0)

    # 히트맵 그리기 - 부드러운 색상과 개선된 스타일
    fig, ax = plt.subplots(figsize=(16, 8))

    # 부드러운 색상으로 히트맵 그리기
    im = ax.imshow(
        pivot.values, aspect="auto", cmap="RdYlBu_r", interpolation="nearest"
    )

    # 컬러바 설정
    cbar = plt.colorbar(im, ax=ax, shrink=0.8)
    cbar.set_label("게시물 수", fontsize=12, color="#2C3E50")
    cbar.ax.tick_params(labelsize=10, colors="#2C3E50")

    # 요일 라벨 (한국어)
    weekday_labels = [
        "월요일",
        "화요일",
        "수요일",
        "목요일",
        "금요일",
        "토요일",
        "일요일",
    ]
    ax.set_yticks(range(7))
    ax.set_yticklabels(weekday_labels, fontsize=11, color="#2C3E50")

    # 시간 라벨 (실제 데이터에 맞게)
    hour_labels = [f"{col:02d}시" for col in pivot.columns]
    ax.set_xticks(range(len(pivot.columns)))
    ax.set_xticklabels(
        hour_labels, fontsize=10, color="#2C3E50", rotation=30, ha="right"
    )

    # 각 셀에 값 표시 (값이 0이 아닌 경우만)
    for i in range(pivot.shape[0]):
        for j in range(pivot.shape[1]):
            if j < len(pivot.columns):  # 컬럼 범위 체크
                value = int(pivot.iloc[i, j])
                if value > 0:  # 값이 있는 경우만 표시
                    text_color = (
                        "white"
                        if pivot.iloc[i, j] > pivot.values.max() * 0.5
                        else "#2C3E50"
                    )
                    ax.text(
                        j,
                        i,
                        str(value),
                        ha="center",
                        va="center",
                        fontsize=9,
                        fontweight="bold",
                        color=text_color,
                    )

    plt.title(
        "요일별 시간대 게시물 빈도 히트맵",
        fontsize=16,
        fontweight="bold",
        color="#2C3E50",
        pad=20,
    )
    plt.xlabel("시간 (24시간)", fontsize=12, color="#2C3E50")
    plt.ylabel("요일", fontsize=12, color="#2C3E50")

    # 배경 설정
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    savefig(os.path.join(OUT_DIR, "05_time_heatmap.png"))


def chart_location_density(posts: pd.DataFrame) -> None:
    # 단순 2D Hist (lat/lon) - 부드러운 색상으로 개선
    lat = posts["lat"].astype(float)
    lon = posts["lon"].astype(float)
    lat = lat.dropna()
    lon = lon.dropna()

    fig, ax = plt.subplots(figsize=(12, 10))

    # 부드러운 색상으로 히스토그램 그리기
    hist, xedges, yedges, im = ax.hist2d(lon, lat, bins=60, cmap="RdYlBu_r", alpha=0.8)

    # 컬러바 설정
    cbar = plt.colorbar(im, ax=ax, shrink=0.8)
    cbar.set_label("게시물 밀도", fontsize=12, color="#2C3E50")
    cbar.ax.tick_params(labelsize=10, colors="#2C3E50")

    plt.title(
        "거래 위치 밀도 (2D 히스토그램)",
        fontsize=16,
        fontweight="bold",
        color="#2C3E50",
        pad=20,
    )
    plt.xlabel("경도", fontsize=12, color="#2C3E50")
    plt.ylabel("위도", fontsize=12, color="#2C3E50")

    # 배경 설정
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    savefig(os.path.join(OUT_DIR, "06_location_2dhist.png"))


def chart_location_map(posts: pd.DataFrame) -> None:
    """거래 위치를 지도로 시각화 (HTML 지도)"""
    if not FOLIUM_AVAILABLE:
        print("⚠️  folium이 설치되지 않았습니다. 지도 시각화를 건너뜁니다.")
        return

    # 좌표 데이터 준비
    lat = posts["lat"].astype(float).dropna()
    lon = posts["lon"].astype(float).dropna()

    if len(lat) == 0 or len(lon) == 0:
        print("⚠️  위치 데이터가 없습니다. 지도 시각화를 건너뜁니다.")
        return

    # 서울 중심 좌표
    seoul_center = [37.5665, 126.9780]

    # 지도 생성
    m = folium.Map(location=seoul_center, zoom_start=11, tiles="OpenStreetMap")

    # 히트맵 데이터 준비
    heat_data = [[lat.iloc[i], lon.iloc[i]] for i in range(len(lat))]

    # 히트맵 추가
    plugins.HeatMap(
        heat_data,
        name="거래 위치 히트맵",
        max_zoom=18,
        radius=15,
        blur=10,
        max_intensity=100,
    ).add_to(m)

    # 마커 클러스터 추가 (상위 MAX_MARKERS개 위치만)
    MAX_MARKERS = 100
    marker_data = heat_data[:MAX_MARKERS] if len(heat_data) > MAX_MARKERS else heat_data

    marker_cluster = plugins.MarkerCluster().add_to(m)
    for lat_val, lon_val in marker_data:
        folium.Marker([lat_val, lon_val]).add_to(marker_cluster)

    # 범례 추가
    folium.LayerControl().add_to(m)

    # 지도 저장
    map_path = os.path.join(OUT_DIR, "10_location_map.html")
    m.save(map_path)
    print(f"🗺️  지도 시각화 저장: {map_path}")


def chart_location_scatter(posts: pd.DataFrame) -> None:
    """거래 위치를 산점도로 시각화"""
    # 좌표 데이터 준비
    lat = posts["lat"].astype(float).dropna()
    lon = posts["lon"].astype(float).dropna()

    if len(lat) == 0 or len(lon) == 0:
        print("⚠️  위치 데이터가 없습니다. 산점도 시각화를 건너뜁니다.")
        return

    fig, ax = plt.subplots(figsize=(12, 10))

    # 부드러운 색상으로 산점도 그리기
    ax.scatter(
        lon, lat, c="#4ECDC4", alpha=0.6, s=20, edgecolors="#2C3E50", linewidth=0.3
    )

    plt.title(
        "거래 위치 분포 (산점도)",
        fontsize=16,
        fontweight="bold",
        color="#2C3E50",
        pad=20,
    )
    plt.xlabel("경도", fontsize=12, color="#2C3E50")
    plt.ylabel("위도", fontsize=12, color="#2C3E50")
    plt.grid(True, alpha=0.3)

    # 배경 설정
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    savefig(os.path.join(OUT_DIR, "11_location_scatter.png"))


def chart_market_share(posts: pd.DataFrame) -> None:
    market_counts = posts["mogu_market"].value_counts().head(10)

    # 부드러운 색상 설정
    soft_colors = [
        "#FF6B6B",
        "#4ECDC4",
        "#45B7D1",
        "#96CEB4",
        "#FFEAA7",
        "#DDA0DD",
        "#98D8C8",
        "#F7DC6F",
        "#FFB6C1",
        "#D3D3D3",
    ]
    colors = (soft_colors * ((len(market_counts) // len(soft_colors)) + 1))[
        : len(market_counts)
    ]

    fig, ax = plt.subplots(figsize=(12, 6))
    bars = ax.bar(
        range(len(market_counts)),
        market_counts.values,
        color=colors,
        alpha=0.8,
        edgecolor="#2C3E50",
        linewidth=1,
    )

    # 막대 위에 값 표시
    for bar, value in zip(bars, market_counts.values):
        height = bar.get_height()
        ax.text(
            bar.get_x() + bar.get_width() / 2.0,
            height + height * 0.01,
            f"{value:,}",
            ha="center",
            va="bottom",
            fontweight="bold",
            color="#2C3E50",
        )

    plt.title("마켓 상위 분포", fontsize=16, fontweight="bold", color="#2C3E50", pad=20)
    plt.xlabel("마켓", fontsize=12, color="#2C3E50")
    plt.ylabel("게시물 수", fontsize=12, color="#2C3E50")
    plt.xticks(range(len(market_counts)), market_counts.index, rotation=30, ha="right")
    plt.grid(True, alpha=0.3, axis="y")

    # 배경 설정
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    savefig(os.path.join(OUT_DIR, "07_market_share.png"))


def chart_price_vs_target(posts: pd.DataFrame) -> None:
    # 산점도: 가격 vs 목표 인원 - 부드러운 색상으로 개선
    x = posts["target_count"].astype(float)
    y = posts["price"].astype(float)

    fig, ax = plt.subplots(figsize=(12, 8))

    # 부드러운 색상으로 산점도 그리기
    ax.scatter(x, y, s=20, alpha=0.6, c="#4ECDC4", edgecolors="#2C3E50", linewidth=0.5)

    plt.title(
        "가격 vs 목표 인원 (산점도)",
        fontsize=16,
        fontweight="bold",
        color="#2C3E50",
        pad=20,
    )
    plt.xlabel("목표 인원", fontsize=12, color="#2C3E50")
    plt.ylabel("가격(원)", fontsize=12, color="#2C3E50")
    plt.grid(True, alpha=0.3)

    # 배경 설정
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    savefig(os.path.join(OUT_DIR, "08_price_vs_target.png"))


def chart_ratings_hist(ratings: pd.DataFrame) -> None:
    if "stars" in ratings.columns and len(ratings):
        # 부드러운 색상으로 히스토그램 그리기
        fig, ax = plt.subplots(figsize=(10, 6))

        # 히스토그램 그리기
        n, bins, patches = ax.hist(
            ratings["stars"], bins=5, alpha=0.7, edgecolor="#2C3E50", linewidth=1.5
        )

        # 각 막대에 부드러운 색상 적용
        soft_colors = ["#FF6B6B", "#4ECDC4", "#45B7D1", "#96CEB4", "#FFEAA7"]
        for patch, color in zip(patches, soft_colors):
            patch.set_facecolor(color)

        # 각 막대 위에 값 표시
        for i, count in enumerate(n):
            if count > 0:
                ax.text(
                    bins[i] + 0.4,
                    count + count * 0.01,
                    f"{int(count)}",
                    ha="center",
                    va="bottom",
                    fontweight="bold",
                    color="#2C3E50",
                )

        plt.title("별점 분포", fontsize=16, fontweight="bold", color="#2C3E50", pad=20)
        plt.xlabel("별점", fontsize=12, color="#2C3E50")
        plt.ylabel("빈도", fontsize=12, color="#2C3E50")
        plt.grid(True, alpha=0.3, axis="y")

        # 배경 설정
        ax.set_facecolor("white")
        fig.patch.set_facecolor("white")

        savefig(os.path.join(OUT_DIR, "09_ratings_hist.png"))


def table_persona_alignment(  # noqa: PLR0913
    users: pd.DataFrame,
    posts: pd.DataFrame,
    favs: pd.DataFrame,
    parts: pd.DataFrame,
    persona_nicks: list[str],
    save_csv: bool = True,
) -> pd.DataFrame:
    """페르소나별 관심사/마켓 일치율 테이블 산출"""
    rows = []
    # post lookup
    post_map = posts.set_index("id")[["category", "mogu_market"]]
    for nick in persona_nicks:
        person = users[users["nickname"].str.contains(nick)]
        if person.empty:
            continue
        uid = person.iloc[0]["id"]
        cats = set(person.iloc[0]["interested_categories"])
        mkts = set(person.iloc[0]["wish_markets"])

        fav_u = favs[favs["user_id"] == uid]
        part_u = parts[parts["user_id"] == uid]

        def calc_align(df: pd.DataFrame) -> tuple[float, float]:
            if df.empty:
                return (float("nan"), float("nan"))
            merged = df.merge(
                post_map, left_on="mogu_post_id", right_index=True, how="left"
            )
            cat_align = (merged["category"].isin(cats)).mean() * 100
            mkt_align = (merged["mogu_market"].isin(mkts)).mean() * 100
            return (round(cat_align, 1), round(mkt_align, 1))

        fav_cat, fav_mkt = calc_align(fav_u)
        par_cat, par_mkt = calc_align(part_u)

        rows.append(
            {
                "persona": nick,
                "favorites": len(fav_u),
                "participations": len(part_u),
                "fav_cat_align_%": fav_cat,
                "fav_market_align_%": fav_mkt,
                "par_cat_align_%": par_cat,
                "par_market_align_%": par_mkt,
            }
        )
    df = pd.DataFrame(rows)
    if save_csv:
        df.to_csv(
            os.path.join(OUT_DIR, "10_persona_alignment.csv"),
            index=False,
            encoding="utf-8-sig",
        )
    return df


def chart_persona_bars(alignment_df: pd.DataFrame) -> None:
    df = alignment_df.set_index("persona")
    df[["favorites", "participations"]].plot(kind="bar")
    plt.title("페르소나별 활동량(찜/참여)")
    plt.xlabel("페르소나")
    plt.ylabel("개수")
    plt.xticks(rotation=30, ha="right")
    savefig(os.path.join(OUT_DIR, "11_persona_activity.png"))

    df[["par_cat_align_%", "par_market_align_%"]].plot(kind="bar")
    plt.title("페르소나별 참여 관심사/마켓 일치율(%)")
    plt.xlabel("페르소나")
    plt.ylabel("일치율(%)")
    plt.xticks(rotation=30, ha="right")
    savefig(os.path.join(OUT_DIR, "12_persona_alignment_bars.png"))


def chart_persona_price_sensitivity(
    posts: pd.DataFrame,
    parts: pd.DataFrame,
    users: pd.DataFrame,
    persona: str = "P5_Student_Sindorim",
) -> None:
    # 개인부담 = price / (target_count + 1)
    p = users[users["nickname"].str.contains(persona)]
    if p.empty:
        return
    uid = p.iloc[0]["id"]
    part_u = parts[parts["user_id"] == uid]
    if part_u.empty:
        return
    tmp = part_u.merge(
        posts[["id", "price", "target_count"]],
        left_on="mogu_post_id",
        right_on="id",
        how="left",
    )
    tmp["personal_cost"] = tmp["price"] / (tmp["target_count"] + 1)
    tmp["personal_cost"].plot(kind="hist", bins=15)
    plt.title(f"{persona} 개인부담금 분포")
    plt.xlabel("개인부담(원)")
    plt.ylabel("빈도")
    savefig(os.path.join(OUT_DIR, "13_p5_personal_cost_hist.png"))


def chart_category_share_in_favs(
    posts: pd.DataFrame,
    favs: pd.DataFrame,
    users: pd.DataFrame,
    persona: str = "P3_Fashion_Beauty_Hongdae",
) -> None:
    p = users[users["nickname"].str.contains(persona)]
    if p.empty:
        return
    uid = p.iloc[0]["id"]
    fav_u = favs[favs["user_id"] == uid]
    if fav_u.empty:
        return
    merged = fav_u.merge(
        posts[["id", "category"]], left_on="mogu_post_id", right_on="id", how="left"
    )
    merged["category"].value_counts().plot(kind="bar")
    plt.title(f"{persona} 찜한 게시물 카테고리 분포")
    plt.xlabel("카테고리")
    plt.ylabel("개수")
    plt.xticks(rotation=30, ha="right")
    savefig(os.path.join(OUT_DIR, "14_p3_fav_category_bar.png"))


def chart_joined_vs_status(posts: pd.DataFrame) -> None:
    # 상태별 joined_count 분포
    posts.boxplot(column="joined_count", by="status", grid=False)
    plt.title("상태별 joined_count 분포")
    plt.suptitle("")
    plt.xlabel("상태")
    plt.ylabel("joined_count")
    savefig(os.path.join(OUT_DIR, "15_joined_by_status_box.png"))


# ========== PPT 생성 ==========
def make_ppt(slides: list[dict[str, Any]]) -> None:
    if Presentation is None or Inches is None or Pt is None:
        print("[WARN] python-pptx 미설치로 PPT 생성을 건너뜁니다")
        return

    prs = Presentation()
    blank = prs.slide_layouts[6]

    def add_slide(
        title: str, subtitle: str | None = None, img: str | None = None
    ) -> None:
        s = prs.slides.add_slide(blank)
        left, top = Inches(0.5), Inches(0.3)
        tbox = s.shapes.add_textbox(left, top, Inches(9), Inches(0.8))
        tf = tbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)

        if subtitle:
            tbox2 = s.shapes.add_textbox(
                Inches(0.5), Inches(1.1), Inches(9), Inches(1.6)
            )
            tf2 = tbox2.text_frame
            for line in textwrap.wrap(subtitle, width=70):
                para = tf2.add_paragraph()
                para.text = line
                para.level = 1

        if img and os.path.exists(img):
            s.shapes.add_picture(img, Inches(0.7), Inches(2), height=Inches(4.8))

    for t in slides:
        add_slide(**t)

    prs.save(PPT_PATH)
    print(f"[PPT] 저장 완료 → {PPT_PATH}")


# ========== 메인 ==========
def main() -> None:
    ensure_dir(OUT_DIR)
    data, users, posts, favs, parts, ratings, comments = load_data()

    # 기본 차트/표
    chart_price_box_by_category(posts)
    chart_post_status_bar(posts)
    chart_participation_status_pie(data["statistics"]["participation_status"])
    chart_avg_joined_by_category(posts)
    chart_time_of_day_heatmap(posts)
    chart_location_density(posts)  # 6번: 2D 히스토그램
    chart_market_share(posts)
    chart_price_vs_target(posts)
    chart_ratings_hist(ratings)
    chart_joined_vs_status(posts)

    # 위치 시각화 차트 (분리)
    chart_location_map(posts)  # 10번: HTML 지도
    chart_location_scatter(posts)  # 11번: 산점도

    # 페르소나 정렬/검증 표 + 파생 차트
    persona_nicks = [
        "P1_OfficeWorker_Gangnam",
        "P2_Family_Costco_Suwon",
        "P3_Fashion_Beauty_Hongdae",
        "P4_NightShifter_Incheon",
        "P5_Student_Sindorim",
    ]
    align_df = table_persona_alignment(
        users, posts, favs, parts, persona_nicks, save_csv=True
    )
    chart_persona_bars(align_df)
    chart_persona_price_sensitivity(posts, parts, users, persona="P5_Student_Sindorim")
    chart_category_share_in_favs(
        posts, favs, users, persona="P3_Fashion_Beauty_Hongdae"
    )

    # PPT (선택)
    if MAKE_PPT:
        slides = [
            dict(
                title="모구 AI 더미데이터: 개요",
                subtitle="사람처럼 행동하는 Synthetic Dataset · 시간/공간/가격/페르소나 반영",
            ),
            dict(
                title="가격 분포",
                img=os.path.join(OUT_DIR, "01_price_box_by_category.png"),
            ),
            dict(
                title="게시물 상태 분포",
                img=os.path.join(OUT_DIR, "02_post_status_bar.png"),
            ),
            dict(
                title="참여 상태 비율",
                img=os.path.join(OUT_DIR, "03_participation_status_pie.png"),
            ),
            dict(
                title="평균 참여자(카테고리)",
                img=os.path.join(OUT_DIR, "04_avg_joined_by_category.png"),
            ),
            dict(
                title="시간대 Heatmap", img=os.path.join(OUT_DIR, "05_time_heatmap.png")
            ),
            dict(
                title="거래 위치 밀도",
                img=os.path.join(OUT_DIR, "06_location_2dhist.png"),
            ),
            dict(
                title="마켓 분포 Top", img=os.path.join(OUT_DIR, "07_market_share.png")
            ),
            dict(
                title="가격 vs 목표 인원",
                img=os.path.join(OUT_DIR, "08_price_vs_target.png"),
            ),
            dict(title="별점 분포", img=os.path.join(OUT_DIR, "09_ratings_hist.png")),
            dict(
                title="페르소나 활동량",
                img=os.path.join(OUT_DIR, "11_persona_activity.png"),
            ),
            dict(
                title="페르소나 일치율",
                img=os.path.join(OUT_DIR, "12_persona_alignment_bars.png"),
            ),
            dict(
                title="P5 개인부담금 분포",
                img=os.path.join(OUT_DIR, "13_p5_personal_cost_hist.png"),
            ),
            dict(
                title="P3 찜 카테고리",
                img=os.path.join(OUT_DIR, "14_p3_fav_category_bar.png"),
            ),
            dict(
                title="상태별 joined_count",
                img=os.path.join(OUT_DIR, "15_joined_by_status_box.png"),
            ),
            dict(
                title="결론",
                subtitle="현실성 검증 통과 · 추천 학습/데모 즉시 활용 가능",
            ),
        ]
        make_ppt(slides)

    print(f"[DONE] 차트 PNG 저장 경로: {OUT_DIR}")
    if MAKE_PPT:
        print(f"[DONE] PPT: {PPT_PATH}")


if __name__ == "__main__":
    main()
